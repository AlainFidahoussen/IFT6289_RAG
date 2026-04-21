#!/usr/bin/env python3
"""Generate poster.pptx from the IFT6289 ViDoRe v3 poster content."""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Pt

# ── Colors ────────────────────────────────────────────────────────────────────
BLUE    = RGBColor(0x1b, 0x3a, 0x6b)
ORANGE  = RGBColor(0xc9, 0x7b, 0x2c)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1d, 0x1d, 0x1d)
GREY    = RGBColor(0x66, 0x66, 0x66)
BG      = RGBColor(0xfa, 0xfa, 0xf7)
WIN_BG  = RGBColor(0xe6, 0xf2, 0xe6)
WIN_FG  = RGBColor(0x16, 0x65, 0x34)
LOSE_BG = RGBColor(0xfb, 0xe6, 0xe6)
LOSE_FG = RGBColor(0x99, 0x1b, 0x1b)
HEAD_BG = RGBColor(0xe9, 0xec, 0xf2)
CALL_BG = RGBColor(0xfd, 0xf5, 0xe7)

# ── Slide setup ───────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Mm(1189)
prs.slide_height = Mm(654)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG

# ── Layout ────────────────────────────────────────────────────────────────────
PX     = Mm(22)
PY     = Mm(22)
CG     = Mm(14)
W      = Mm(1189)
H      = Mm(654)
HDR_H  = Mm(55)
COL_W  = int((W - 2 * PX - 3 * CG) / 4)
MAIN_Y = PY + HDR_H + CG
CX     = [int(PX + i * (COL_W + CG)) for i in range(4)]

# ── Helpers ───────────────────────────────────────────────────────────────────

def rect(x: int, y: int, w: int, h: int, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.width = 0


def txb(x: int, y: int, w: int, h: int):
    shape = slide.shapes.add_textbox(x, y, w, h)
    shape.text_frame.word_wrap = True
    return shape.text_frame


def add_run(p, text: str, size: int, bold=False, italic=False, color: RGBColor = DARK):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def h2(x: int, y: int, text: str) -> int:
    h = Mm(14)
    rect(x, y, COL_W, h, BLUE)
    tf = txb(x + Mm(3), y + Mm(1.5), COL_W - Mm(4), h)
    p = tf.paragraphs[0]
    add_run(p, text, 19, bold=True, color=WHITE)
    return y + h + Mm(4)


def h3(x: int, y: int, text: str) -> int:
    rect(x, y + Mm(1.5), Mm(2), Mm(9), ORANGE)
    tf = txb(x + Mm(5), y, COL_W - Mm(5), Mm(12))
    p = tf.paragraphs[0]
    add_run(p, text, 15, bold=True, color=BLUE)
    return y + Mm(13)


def body(x: int, y: int, text: str, size=13, color: RGBColor = DARK, h_est=None) -> int:
    est = h_est or Mm(max(10, len(text) // 5))
    tf = txb(x, y, COL_W, est)
    p = tf.paragraphs[0]
    add_run(p, text, size, color=color)
    return y + est + Mm(2)


def bullets(x: int, y: int, items: list[tuple[str, str]], size=12) -> int:
    row_h = Mm(size * 0.72)
    for bold_part, rest in items:
        tf = txb(x + Mm(2), y, COL_W - Mm(2), row_h + Mm(5))
        p = tf.paragraphs[0]
        add_run(p, "• ", size, color=ORANGE)
        if bold_part:
            add_run(p, bold_part + ": ", size, bold=True)
        if rest:
            add_run(p, rest, size)
        y += row_h + Mm(3)
    return y + Mm(2)


def callout_box(x: int, y: int, text: str, size=12) -> int:
    est = Mm(max(14, len(text) // 4))
    rect(x, y, COL_W, est, CALL_BG)
    rect(x, y, Mm(2.5), est, ORANGE)
    tf = txb(x + Mm(5), y + Mm(3), COL_W - Mm(7), est - Mm(5))
    p = tf.paragraphs[0]
    add_run(p, text, size)
    return y + est + Mm(4)


def table(x: int, y: int, data: list[list[str]], col_fracs: list[float],
          wins: set[tuple[int, int]], loses: set[tuple[int, int]]) -> int:
    nrows, ncols = len(data), len(data[0])
    row_h = Mm(10)
    tbl_h = nrows * row_h
    tbl = slide.shapes.add_table(nrows, ncols, x, y, COL_W, tbl_h).table
    for c, frac in enumerate(col_fracs):
        tbl.columns[c].width = int(COL_W * frac)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.bold = r == 0 or (r == nrows - 1 and c == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEAD_BG
                p.font.color.rgb = BLUE
            elif (r, c) in wins:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WIN_BG
                p.font.color.rgb = WIN_FG
                p.font.bold = True
            elif (r, c) in loses:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LOSE_BG
                p.font.color.rgb = LOSE_FG
                p.font.bold = True
    return y + tbl_h + Mm(5)


# ── HEADER ────────────────────────────────────────────────────────────────────
tf = txb(PX, PY, W - 2 * PX, HDR_H)

p0 = tf.paragraphs[0]
p0.alignment = PP_ALIGN.CENTER
add_run(p0, "Does adding image retrieval help a strong text-based RAG?", 44, bold=True, color=BLUE)

p1 = tf.add_paragraph()
p1.alignment = PP_ALIGN.CENTER
add_run(p1, "Reproducing ", 23, color=GREY)
add_run(p1, "ViDoRe v3", 23, color=ORANGE)
add_run(p1, " with open-weights models", 23, color=GREY)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
add_run(p2, "Alain Fidahoussen · Munyeong Kim · Aftab Gazali", 14, color=DARK)

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
add_run(p3, "Université de Montréal · MILA", 12, italic=True, color=GREY)

rect(PX, PY + HDR_H, W - 2 * PX, Mm(0.5), DARK)

# ── COLUMN 1: WHY HYBRID RAG ──────────────────────────────────────────────────
x, y = CX[0], MAIN_Y

y = h2(x, y, "Why hybrid RAG?")
y = body(x, y, (
    "RAG pipelines over visually rich PDFs face a basic tension: the answer may live "
    "in text (paragraphs, tables) or in visual elements (charts, figures, handwriting) "
    "that text extractors drop or flatten."
), h_est=Mm(30))
y = bullets(x, y, [
    ("Textual retrieval", "parser → markdown → dense embedding. Precise on text; blind to figures."),
    ("Visual retrieval", "page-image encoder (ColPali / ColEmbed). Sees layout and figures; less precise on fine-grained text."),
])
y = callout_box(x, y, (
    "Hybrid concatenates the top-k of both streams and feeds them to a multimodal generator. "
    "The intuition: the two streams capture complementary evidence. Does that hold in practice?"
))

y = h2(x, y, "Visual retrieval")
y = body(x, y, (
    "Each PDF page is encoded as pixels by a vision-language model, once, offline. "
    "No OCR, no markdown, no chunking. A page is its own atom."
), h_est=Mm(22))

y = h3(x, y, "ColPali (Faysse et al., 2025)")
y = body(x, y, (
    "Late interaction: each page is encoded into one embedding per image patch, each query into one per token. "
    "Scoring uses MaxSim: for every query token, find its closest page patch and sum those best-match scores. "
    'A query "what is the Q3 revenue?" can match "Q3" to the right table column and "revenue" to the right row.'
), h_est=Mm(38))

y = h3(x, y, "ColEmbed (Xu et al., 2025) — our visual stream")
y = body(x, y, (
    "Same late-interaction architecture, two upgrades: stronger backbone Llama-Nemotron "
    "(NVIDIA's document-focused VLM) and broader synthetic query-page training with hard-negative mining."
), h_est=Mm(28))

y = h2(x, y, "The ViDoRe v3 benchmark")
y = body(x, y, (
    "Loison et al. (2026): 10 corpora, 26 000 pages, 3 099 human-verified queries in 6 languages. "
    "We evaluate on 5 public subsets (3 EN + 2 FR): 10 673 pages and 2 132 queries."
), h_est=Mm(28))

y = h3(x, y, "What makes it hard")
y = bullets(x, y, [
    ("Visually rich documents", "annual reports, scientific books, pharma regulations. Tables, charts, figures carry real evidence."),
    ("7 query types", "from open-ended synthesis to single-number lookup."),
    ("Cross-lingual", "queries in 6 languages paired with EN or FR source documents."),
    ("Decoupled evaluation", "best retriever is not always the best generator (NDCG@10 vs pass@1)."),
])

# ── COLUMN 2: METHOD ──────────────────────────────────────────────────────────
x, y = CX[1], MAIN_Y

y = h2(x, y, "Method")
y = callout_box(x, y, (
    "Question: with the text stream fixed at DeepSeek-OCR-2 + Jina v4 + zerank-2, "
    "does adding a visual retrieval stream (ColEmbed) improve end-to-end answer accuracy?"
))

y = h3(x, y, "Pipeline")
y = bullets(x, y, [
    ("Text stream", "PDF → DeepSeek-OCR-2 → markdown → Jina v4 (dense vectors) → zerank-2 (reranker) → Top-5 text pages"),
    ("Image stream", "PDF page images → ColEmbed (image encoder) → Top-5 image pages"),
    ("Fusion + eval", "Top-5 text + Top-5 image → qwen3.5:35b (generator) → Answer → llama3.1:8b (judge) → pass@1"),
])

y = h3(x, y, "Model choices")
y = bullets(x, y, [
    ("DeepSeek-OCR-2", "open-weights document OCR with layout-preserving markdown (HTML tables, LaTeX)"),
    ("qwen3.5:35b", "multimodal generator, Q4-quantized MoE, fits on one 48 GB GPU, 262k context"),
    ("llama3.1:8b", "judge from a different model family (Meta vs Alibaba) to reduce self-preference bias"),
])
y = body(x, y, "All components are open-weights and run on a single GPU.", 11, GREY, h_est=Mm(12))

y = h3(x, y, "DeepSeek-OCR-2")
y = body(x, y, (
    "Instead of consuming visual tokens top-left to bottom-right, learnable causal flow queries "
    "re-order them to a more logical reading path, preserving document structure in the output markdown."
), h_est=Mm(26))

pic = slide.shapes.add_picture("poster/assets/deepseek_ocr2_fig3.jpg", x, y, width=COL_W)
y += pic.height + Mm(5)

y = h3(x, y, "Experimental conditions (7 total)")
y = bullets(x, y, [
    ("jina_nemo / jina_deepseek", "text only, top-5, no reranker"),
    ("jina_nemo_reranked / jina_deepseek_reranked", "text only, top-5, zerank-2 reranker"),
    ("colembed", "image only, top-5"),
    ("hybrid_nemo / hybrid_deepseek", "image top-5 + text top-5, zerank-2 on text stream"),
])

# ── COLUMN 3: RESULTS ─────────────────────────────────────────────────────────
x, y = CX[2], MAIN_Y

y = h2(x, y, "Results")
y = h3(x, y, "Retrieval quality (NDCG@10)")

y = table(x, y,
    data=[
        ["Subset",                 "ColEmbed (image)", "Jina + zerank-2 (text)"],
        ["computer_science (EN)",  "78.0",             "82.4"],
        ["finance_en (EN)",        "68.6",             "65.7"],
        ["pharmaceuticals (EN)",   "67.2",             "65.1"],
        ["physics (FR)",           "47.8",             "46.6"],
        ["finance_fr (FR)",        "47.6",             "48.1"],
        ["avg",                    "61.8",             "61.6"],
    ],
    col_fracs=[0.50, 0.25, 0.25],
    wins={(1, 2), (2, 1), (3, 1), (4, 1), (5, 2)},
    loses=set(),
)

y = h3(x, y, "End-to-end accuracy (pass@1 %)")

y = table(x, y,
    data=[
        ["Subset",          "Closed", "Image", "Text",  "Hybrid"],
        ["computer_science","91.6",   "93.0",  "95.3",  "95.3"],
        ["finance_en",      "56.0",   "74.8",  "79.0",  "81.6"],
        ["pharmaceuticals", "68.4",   "83.2",  "88.5",  "89.3"],
        ["physics",         "87.8",   "83.1",  "90.1",  "89.7"],
        ["finance_fr",      "40.6",   "70.3",  "80.3",  "77.2"],
        ["avg",             "68.9",   "80.9",  "86.6",  "86.6"],
    ],
    col_fracs=[0.32, 0.17, 0.17, 0.17, 0.17],
    wins={(1, 3), (1, 4), (2, 4), (3, 4), (4, 3), (5, 3), (6, 4)},
    loses=set(),
)

y = h3(x, y, "pass@1 by query type (all subsets pooled)")

y = table(x, y,
    data=[
        ["Query type",       "Closed", "Image", "Text",  "Hybrid"],
        ["numerical",        "36.4",   "70.4",  "70.7",  "74.1"],
        ["extractive",       "54.3",   "77.4",  "83.1",  "82.6"],
        ["multi-hop",        "68.6",   "84.7",  "84.7",  "91.5"],
        ["enumerative",      "60.3",   "75.8",  "83.5",  "80.9"],
        ["boolean",          "72.1",   "85.0",  "93.2",  "87.8"],
        ["compare-contrast", "71.9",   "80.8",  "84.6",  "85.3"],
        ["open-ended",       "82.8",   "82.6",  "91.9",  "91.5"],
    ],
    col_fracs=[0.32, 0.17, 0.17, 0.17, 0.17],
    wins={(1, 4), (2, 3), (3, 4), (4, 3), (5, 3), (6, 4), (7, 3)},
    loses={(5, 4)},
)

# ── COLUMN 4: TAKEAWAYS + LIMITATIONS + FUTURE WORK ──────────────────────────
x, y = CX[3], MAIN_Y

y = h2(x, y, "Takeaways")
y = callout_box(x, y, "Image retrieval is query-type conditional. Overall, hybrid and text-only are tied at 86.6% pass@1.")

y = bullets(x, y, [
    ("Multi-hop", "+6.8 pts over text: two streams retrieve different pages, giving the generator broader coverage."),
    ("Numerical", "+3.4 pts: visual table layout preserves row/column structure that OCR partly loses."),
    ("Boolean", "-5.4 pts: text already retrieves the fact; extra image pages dilute the prompt."),
    ("Strong text baseline", "DeepSeek + Jina + zerank-2 matches hybrid at half the retrieval cost (5 vs 10 pages)."),
    ("Retrieval value", "+37.7 pts on numerical, +28.3 on extractive, only +8.7 on open-ended over closed-book."),
    ("Generator bottleneck", "ViDoRe v3 reports +2.6 pts with Gemini 3 Pro. Our null result is likely generator-limited."),
])

y = h2(x, y, "Limitations")
y = h3(x, y, "LLM choices")
y = bullets(x, y, [
    ("Generator qwen3.5:35b", "vision stack weaker than frontier closed models (Gemini 3 Pro, GPT-5.2)."),
    ("Judge llama3.1:8b", "verbosity bias: longer answers get rewarded. Hybrid answers tend to be longer."),
])

y = h3(x, y, "Benchmark and evaluation design")
y = bullets(x, y, [
    ("Binary pass@1", "near-miss and complete miss treated identically."),
    ("No human spot-check", "all judgments from LLM judge, no human validation."),
    ("Fixed fusion", "always top-5+top-5, no deduplication; same page can appear twice."),
    ("Domain-specific", "CS, finance, pharma only. May not transfer to general-purpose RAG."),
])

y = h2(x, y, "Future work")
y = bullets(x, y, [
    ("Stronger generators", "Gemini 3 Pro, Qwen3-VL-235B to isolate the generator-capability effect."),
    ("Verbosity control", "Controlled-verbosity prompt and judge swap to rule out verbosity bias."),
    ("Adaptive fusion", "Weighted combination, deduplication, query-type routing."),
    ("Figure captioning", "Close the figure-placeholder blind spot in the text stream."),
])

# ── Footer ────────────────────────────────────────────────────────────────────
rect(PX, H - PY - Mm(1), W - 2 * PX, Mm(0.5), RGBColor(0xcc, 0xcc, 0xcc))
tf = txb(PX, H - PY, W - 2 * PX, Mm(8))
p = tf.paragraphs[0]
add_run(p, (
    "Text parser: DeepSeek-OCR-2  ·  Text retriever: Jina v4  ·  Reranker: zerank-2  ·  "
    "Visual retriever: ColEmbed (Llama-Nemotron 3B v2)  ·  Generator: qwen3.5:35b  ·  "
    "Judge: llama3.1:8b  ·  Benchmark: ViDoRe v3 (Loison et al., 2026)"
), 9, color=GREY)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "poster/poster.pptx"
prs.save(OUT)
print(f"Saved {OUT}")
